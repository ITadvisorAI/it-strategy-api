import os
import json
import traceback
import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# === Google Drive Setup ===
drive_service = None
try:
    creds_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if creds_json:
        creds = service_account.Credentials.from_service_account_info(
            json.loads(creds_json),
            scopes=["https://www.googleapis.com/auth/drive"]
        )
        drive_service = build("drive", "v3", credentials=creds)
except Exception as e:
    print(f"❌ Drive init failed: {e}")
    traceback.print_exc()

def upload_to_drive(file_path, session_id):
    try:
        folder_id = None
        query = f"name='{session_id}' and mimeType='application/vnd.google-apps.folder'"
        results = drive_service.files().list(q=query, fields="files(id)").execute()
        folders = results.get("files", [])
        if folders:
            folder_id = folders[0]["id"]
        else:
            folder = drive_service.files().create(body={
                "name": session_id,
                "mimeType": "application/vnd.google-apps.folder"
            }, fields="id").execute()
            folder_id = folder["id"]

        file_meta = {"name": os.path.basename(file_path), "parents": [folder_id]}
        media = MediaFileUpload(file_path, resumable=True)
        uploaded = drive_service.files().create(body=file_meta, media_body=media, fields="id").execute()
        return f"https://drive.google.com/file/d/{uploaded['id']}/view"
    except Exception as e:
        print(f"❌ Upload failed: {e}")
        traceback.print_exc()
        return None

def extract_target_recommendations(path):
    wb = load_workbook(path)
    sheet = wb.active
    upgrades = []
    for row in sheet.iter_rows(min_row=2, values_only=True):
        platform, tier, status, recommendation = row[2], row[3], row[4], row[5]
        if recommendation:
            upgrades.append(f"{platform} → {recommendation}")
    return list(set(upgrades))

def process_it_strategy(session_id, email, files, folder_path):
    try:
        os.makedirs(folder_path, exist_ok=True)

        downloaded = []
        for f in files:
            path = os.path.join(folder_path, f["file_name"])
            r = requests.get(f["file_url"], timeout=10)
            with open(path, "wb") as fp:
                fp.write(r.content)
            f["local_path"] = path
            downloaded.append(f)

        # Extract recommendations
        hw_recs, sw_recs = [], []
        for f in downloaded:
            if f["file_type"] == "hardware_gap":
                hw_recs = extract_target_recommendations(f["local_path"])
            elif f["file_type"] == "software_gap":
                sw_recs = extract_target_recommendations(f["local_path"])

        # === Generate DOCX Strategy Report ===
        docx_path = os.path.join(folder_path, "IT Infrastructure Upgrade Strategy.docx")
        doc = Document()
        doc.add_heading("IT Infrastructure Upgrade Strategy", 0)
        doc.add_paragraph(f"Session ID: {session_id}")
        doc.add_paragraph("\n1. Introduction\nThis document outlines the proposed target infrastructure state.")
        doc.add_paragraph("\n2. Hardware Upgrade Plan:\n" + "\n".join(hw_recs) or "No hardware upgrades required.")
        doc.add_paragraph("\n3. Software Upgrade Plan:\n" + "\n".join(sw_recs) or "No software upgrades required.")
        doc.add_paragraph("\n4. Considerations:\nCost, geographic location, and performance constraints were considered.")
        doc.add_paragraph("\n5. Summary:\nThe proposed architecture ensures scalability, security, and future-readiness.")
        doc.save(docx_path)

        # === Generate PPTX Executive Summary ===
        pptx_path = os.path.join(folder_path, "IT Infrastructure Upgrade Executive Report.pptx")
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.shapes.title.text = "IT Upgrade Executive Report"
        slide.placeholders[1].text = f"Session: {session_id}"

        def add_slide(title, items):
            s = ppt.slides.add_slide(ppt.slide_layouts[1])
            s.shapes.title.text = title
            body = s.placeholders[1]
            body.text = ""
            for i in items:
                p = body.text_frame.add_paragraph()
                p.text = i

        add_slide("Hardware Upgrades", hw_recs)
        add_slide("Software Upgrades", sw_recs)
        add_slide("Constraints Considered", ["Cost Optimization", "Geographic Accessibility", "Future-readiness"])

        ppt.save(pptx_path)

        # Upload and update file URLs
        docx_url = upload_to_drive(docx_path, session_id)
        pptx_url = upload_to_drive(pptx_path, session_id)

        for f in downloaded:
            f["file_url"] = upload_to_drive(f["local_path"], session_id)

        downloaded.extend([
            {
                "file_name": os.path.basename(docx_path),
                "file_url": docx_url,
                "file_type": "docx_strategy"
            },
            {
                "file_name": os.path.basename(pptx_path),
                "file_url": pptx_url,
                "file_type": "pptx_strategy"
            }
        ])

        # Forward to next GPT (Target GAP Analysis)
        NEXT_GPT_URL = "https://gap-target-api.onrender.com/start_gap_target"
        payload = {
            "session_id": session_id,
            "email": email,
            "gpt_module": "it_strategy",
            "files": downloaded,
            "status": "complete"
        }
        requests.post(NEXT_GPT_URL, json=payload)

    except Exception as e:
        print(f"🔥 Strategy processing failed: {e}")
        traceback.print_exc()
